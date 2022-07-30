import collections.abc
import pptx
import os

ppt=pptx.Presentation()

# 16:9
ppt.slide_height = 6858000    #设置ppt的高度
ppt.slide_width = 12192000    #设置ppt的宽度

picture=os.listdir('.\image')

for file in sorted(picture,key=lambda x:int(x.split('.')[0])):
    slide = ppt.slides.add_slide(ppt.slide_layouts[3])
    slide.shapes.add_picture(f'.\image\{file}',0,0,ppt.slide_width,ppt.slide_height)


ppt.save('人力资源.pptx')